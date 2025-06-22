# -*- coding: utf-8 -*-
"""
Created on Tue Jun 17 19:26:43 2025

@author: Yiming
"""

import pymupdf
from google import genai
from pptx import Presentation
from pptx.util import Pt
import os

'''
GEMINI_API_KEY = ''
PDF_PATH = 'paper.pdf'
num_slides = 5
OUTPUT_PPTX = 'summary.pptx'
'''
import argparse

def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("GEMINI_API_KEY", help="Your Gemini API Key")
    parser.add_argument("PDF_PATH", help="Path to input PDF file")
    parser.add_argument("num_slides", type=int, help="Number of slides to generate")
    return parser.parse_args()


def extract_pdf_text(path, max_pages=10) -> str:
    doc = pymupdf.open(path)
    text = ""
    for page in doc[:max_pages]:
        text += page.get_text()
    return text


def extract_images_from_pdf(pdf_path, output_folder='images', min_width=72, min_height=72)-> list[str]:
    '''
    Only extract (png,jpg) images, not vector images (eps) from pdf
    min_width in points, 72 pts = 1 inch
    '''
    os.makedirs(output_folder, exist_ok=True)
    doc = pymupdf.open(pdf_path)
    image_paths = []

    for page_num, page in enumerate(doc):
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]  # cross-reference number, int
            rects = page.get_image_rects(xref)
            if not rects:
                continue
            # For each on-page occurrence, keep it if big enough
            for rect in rects:
                if rect.width < min_width or rect.height < min_height:
                    continue

                # Extract the pixmap
                pix = pymupdf.Pixmap(doc, xref)
            
                img_path = os.path.join(output_folder, f"page{page_num+1}_img{img_index+1}.jpg")
                pix.save(img_path)
                pix = None
                image_paths.append(img_path)

    return image_paths  # list of all saved image paths


def generate_slide_content(text: str, model='gemini-2.0-flash') -> str:
    prompt = f"""
    You are an expert scientific summarizer.
    Summarize the following academic paper into no more than **{num_slides}** clean PowerPoint slides.
    You are welcome to use less slides. 
    More advanced expert can use less slides and less bullets to cover key points.
    For each slide output exactly:

    Slide 1:
    Title: <one short sentence>
    Bullets:
    - <up to 5 bullets, max 16 words each>
    - <…>
    - <…>

    Mark “[[FIGURE+number or letter]]” wherever a figure from the paper should be placed (we’ll move it later).
    Do not include any other words, symbols, or slide numbers.

    Paper:
    {text[:10000]}  # truncate to avoid token overflow
    """
    response = client.models.generate_content(model=model, contents= prompt)
    return response.text

MAX_WORDS_PER_BULLET=16
MAX_BULLETS = 5

def clean_bullet(line):
    # drop everything after MAX_WORDS_PER_BULLET words
    words = line.strip().split()
    return ' '.join(words[:MAX_WORDS_PER_BULLET])

def create_ppt(slide_text, filename):
    prs = Presentation()
    slides = slide_text.strip().split("Slide")[1:]  # crude splitter, improve in the future
    for i, slide in enumerate(slides[:5], 1):  # force max 5 slides
        # parse
        lines = [l.strip() for l in slide.splitlines() if l.strip()]
        bullets = []
        for l in lines:
            if l.startswith("-"):
                b = l.lstrip("- ").strip()  # bullets are the lines start with '- '
                bullets.append(clean_bullet(b))
            elif l.startswith('Title: '):  # title is the line starts with 'Title: '
                title = l.lstrip('Title: ').strip()
            if len(bullets) >= MAX_BULLETS:
                break
    
        # build slide
        layout = prs.slide_layouts[1]  # title+content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(18)

    prs.save(filename)


if  __name__ == '__main__':
    args = parse_args()
    GEMINI_API_KEY = args.GEMINI_API_KEY
    PDF_PATH = args.PDF_PATH
    num_slides = args.num_slides
    OUTPUT_PPTX = os.path.splitext(os.path.basename(PDF_PATH))[0] + "_summary.pptx"
    
    client = genai.Client(api_key=GEMINI_API_KEY)
    
    pdf_text = extract_pdf_text(PDF_PATH)
    # image_paths = extract_images_from_pdf(PDF_PATH)
    slide_text = generate_slide_content(pdf_text, model= 'gemini-2.0-flash')
    create_ppt(slide_text, OUTPUT_PPTX)
